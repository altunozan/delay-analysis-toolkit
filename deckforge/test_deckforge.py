"""Headless test harness for DeckForge (no Streamlit needed).

Builds every chart type from the bundled samples, renders each through BOTH
renderers, asserts the shared maths, and round-trips the generated .pptx
through python-pptx to prove the deck is valid and populated. Also exercises
deck automation (agenda slides, harmonised scales, grid layouts).

Usage:
    python3 deckforge/test_deckforge.py
"""

from __future__ import annotations

import io
import os
import sys
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import samples
from data_io import (
    frame_to_bar, frame_to_butterfly, frame_to_gantt, frame_to_line,
    frame_to_mekko, frame_to_pie, frame_to_process, frame_to_scatter,
    frame_to_table, frame_to_waterfall,
)
from render_plotly import render as render_web
from render_pptx import build_deck
from specs import DeltaArrow, weekend_ranges
from theme import DEFAULT_THEME


def build_core_specs():
    """The original 6 chart types — full Gantt feature set included."""
    bar = frame_to_bar(
        samples.bar_frame(), title="Revenue by workstream", mode="stacked",
        orientation="vertical", show_values=True, show_totals=True,
        deltas=[DeltaArrow("FY22", "FY25", "absolute"),
                DeltaArrow("FY22", "FY25", "cagr")],
    )
    clustered = frame_to_bar(
        samples.bar_frame(), title="Clustered view", mode="clustered",
        orientation="horizontal", show_values=True, show_totals=False,
        deltas=[],
    )
    wf = frame_to_waterfall(samples.waterfall_frame(),
                            title="Delay build-up (days)",
                            show_values=True, show_connectors=True)
    gantt = frame_to_gantt(
        samples.gantt_frame(), title="Summary programme",
        today=datetime(2026, 7, 7),
        curtains_df=samples.gantt_curtains_frame(),
        datelines_df=samples.gantt_datelines_frame(),
        brackets_df=samples.gantt_brackets_frame(),
        show_date_labels=True, show_durations=True, show_remarks=True,
        weekend_shading=True,
    )
    mekko = frame_to_mekko(samples.mekko_frame(), title="Market mekko",
                           show_values=True)
    table = frame_to_table(samples.table_frame(), title="Status table")
    return [bar, clustered, wf, gantt, mekko, table]


def build_extended_specs():
    """Everything added in the A-sweep: combo, break, sort, units, pie,
    butterfly, scatter/bubble, process flow."""
    combo = frame_to_bar(
        samples.bar_frame(), title="Revenue vs fit-out trend",
        mode="stacked", orientation="vertical", show_values=True,
        show_totals=True, deltas=[], unit="m",
        line_columns=["Fit-out"],
    )
    broken = frame_to_bar(
        samples.bar_frame(), title="Broken axis", mode="stacked",
        orientation="vertical", show_values=True, show_totals=True,
        deltas=[], axis_break=(250.0, 550.0),
    )
    sorted_bar = frame_to_bar(
        samples.bar_frame(), title="Sorted clustered", mode="clustered",
        orientation="vertical", show_values=True, show_totals=False,
        deltas=[], sort="desc",
    )
    line = frame_to_line(samples.line_frame(), title="Progress S-curve",
                         mode="line", show_values=False)
    area = frame_to_line(samples.line_frame(), title="Area view",
                         mode="area")
    pie = frame_to_pie(samples.pie_frame(), title="Cost split")
    doughnut = frame_to_pie(samples.pie_frame(), title="Cost split (ring)",
                            doughnut=True)
    butterfly = frame_to_butterfly(samples.butterfly_frame(),
                                   title="Planned vs actual crew")
    scatter = frame_to_scatter(
        samples.scatter_frame(), title="Package risk map",
        x_title="Delay risk", y_title="Cost impact", quadrants=(60.0, 60.0))
    process = frame_to_process(samples.process_frame(),
                               title="Claim process", highlight=2)
    return [combo, broken, sorted_bar, line, area, pie, doughnut,
            butterfly, scatter, process]


def check_core_maths(specs) -> None:
    bar, _clustered, wf, gantt, mekko, _table = specs

    assert bar.totals() == [690, 793, 908, 1070], bar.totals()
    (i, j, t0, t1, label) = bar.resolved_deltas()[0]
    assert (i, j) == (0, 3) and label == "+380", (i, j, label)
    cagr_label = bar.resolved_deltas()[1][4]
    assert cagr_label.startswith("CAGR +15."), cagr_label

    resolved = wf.resolved()
    assert [round(s.cumulative) for s in resolved] == [45, 107, 125, 100, 100]
    assert resolved[3].kind == "decrease" and resolved[3].base == 100.0
    assert resolved[-1].kind == "total"

    assert gantt.span()[0].year == 2026 and len(gantt.groups()) == 5
    assert len(gantt.items) == 9 and len(gantt.rows()) == 8
    assert any(it.style == "striped" for it in gantt.items)
    assert len(gantt.curtains) == 1 and gantt.curtains[0].label
    assert len(gantt.all_date_lines()) == 2
    lvls = {lvl for _, lvl in gantt.bracket_levels()}
    assert lvls == {0, 1}, lvls
    assert weekend_ranges(*gantt.span()) == []
    assert len(weekend_ranges(datetime(2026, 1, 1),
                              datetime(2026, 2, 1))) >= 4
    assert abs(sum(mekko.column_shares()) - 1.0) < 1e-9
    print("Core spec maths: OK")


def check_extended_maths(specs) -> None:
    (combo, broken, sorted_bar, line, area, pie, doughnut, butterfly,
     scatter, process) = specs

    assert combo.overlay_lines and combo.overlay_lines[0].name == "Fit-out"
    assert combo.unit == "m" and combo.fmt(1_500_000) != "1,500,000"

    t, ticks = broken.axis_break.transform(1200.0)
    assert t(100) < t(250) < t(550) < t(1000)
    gap_span = t(550) - t(250)
    assert gap_span < 0.10, gap_span  # compressed, not linear

    totals = [sum((s.values[i] or 0) for s in sorted_bar.series)
              for i in range(len(sorted_bar.categories))]
    assert totals == sorted(totals, reverse=True), totals

    assert pie.values and abs(sum(pie.values) - 100) < 1e-9
    assert doughnut.doughnut is True

    assert butterfly.left.name and butterfly.right.name
    assert len(butterfly.left.values) == len(butterfly.categories)

    assert scatter.bubble is True and len(scatter.groups()) == 4
    assert scatter.quadrants == (60.0, 60.0)

    assert process.highlight == 2 and len(process.steps) == 6

    from specs import decode_cell
    assert decode_cell("rag:green") == ("●", "positive")
    assert decode_cell("hb:75") == ("◕", "accent")
    assert decode_cell("check") == ("✓", "positive")
    assert decode_cell("plain text") == ("plain text", None)
    print("Extended spec maths: OK")


def main() -> None:
    core = build_core_specs()
    extended = build_extended_specs()
    all_specs = core + extended

    check_core_maths(core)
    check_extended_maths(extended)

    # --- plotly renderer for every spec ---------------------------------
    for spec in all_specs:
        fig = render_web(spec, DEFAULT_THEME)
        assert fig.data or fig.layout.shapes, \
            f"empty plotly figure for {type(spec).__name__}"
        fig.to_json()  # must serialise cleanly
    print(f"Plotly: {len(all_specs)} figures rendered OK")

    # --- pptx: core deck, simple one-chart-per-slide --------------------
    blob = build_deck(core, DEFAULT_THEME, deck_title="Test deck")
    assert len(blob) > 20_000, "suspiciously small pptx"

    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob))
    n_slides = len(prs.slides._sldIdLst)  # noqa: SLF001
    assert n_slides == len(core) + 1, n_slides
    charts = sum(1 for s in prs.slides for sh in s.shapes if sh.has_chart)
    tables = sum(1 for s in prs.slides for sh in s.shapes if sh.has_table)
    shapes = sum(len(s.shapes) for s in prs.slides)
    assert charts == 3, f"expected 3 native charts, got {charts}"
    assert tables == 1, f"expected 1 native table, got {tables}"
    print(f"PPTX (core): {n_slides} slides, {charts} native charts, "
          f"{tables} native table, {shapes} shapes total — round-trip OK")

    out = os.path.join(os.path.dirname(__file__), "sample_output.pptx")
    with open(out, "wb") as fh:
        fh.write(blob)
    print(f"Sample deck written to {out}")

    # --- pptx: extended chart types, native-chart count check -----------
    blob2 = build_deck(extended, DEFAULT_THEME)
    prs2 = Presentation(io.BytesIO(blob2))
    assert len(prs2.slides._sldIdLst) == len(extended)  # noqa: SLF001
    charts2 = sum(1 for s in prs2.slides for sh in s.shapes if sh.has_chart)
    # combo (2: bar+overlay), sorted-clustered, line, area, pie, doughnut,
    # butterfly, scatter = 9 native charts; broken-axis is shapes only.
    assert charts2 == 9, f"expected 9 native charts, got {charts2}"
    print(f"PPTX (extended): {len(prs2.slides._sldIdLst)} slides, "  # noqa: SLF001
          f"{charts2} native charts — round-trip OK")

    # --- deck automation: agenda + harmonised scale + grid layout -------
    grouped = [[core[0]], [core[1], extended[2]], [extended[0]]]
    blob3 = build_deck(
        core[:2] + [extended[0], extended[2]], DEFAULT_THEME,
        deck_title="Automation test", agenda="chapters",
        harmonise_bars=True, groups=grouped)
    prs3 = Presentation(io.BytesIO(blob3))
    # title + 3 chapter-divider/content pairs = 1 + 3*2 = 7
    assert len(prs3.slides._sldIdLst) == 7, len(prs3.slides._sldIdLst)  # noqa: SLF001
    slide_texts = [
        "".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        for s in prs3.slides
    ]
    assert any("Agenda" in t for t in slide_texts)
    print("Deck automation (agenda/harmonise/grid): OK")

    # --- PNG export sanity (kaleido) -------------------------------------
    try:
        png = render_web(core[0], DEFAULT_THEME).to_image(
            format="png", width=400, height=300)
        assert len(png) > 500
        print(f"PNG export: OK ({len(png)} bytes)")
    except Exception as exc:  # noqa: BLE001 — optional dependency
        print(f"PNG export: SKIPPED ({exc})")

    print("\nAll DeckForge assertions passed.")


if __name__ == "__main__":
    main()
