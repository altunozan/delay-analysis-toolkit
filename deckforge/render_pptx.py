"""ChartSpec → native PowerPoint slides (the deck renderer).

Strategy per chart type:
- Bars (stacked/clustered/100%)  → NATIVE editable chart. Totals, difference
  arrows and CAGR arrows are overlay shapes positioned with an explicit
  plot-area geometry model (the value axis max is pinned so positions are
  deterministic).
- Waterfall                      → NATIVE stacked column with an invisible
  base series (the same trick think-cell uses), per-point signed labels and
  shape connectors.
- Gantt & Marimekko              → drawn entirely as shapes (PowerPoint has
  no native chart for these), fully deterministic.
- Table                          → native PowerPoint table.
"""

from __future__ import annotations

import io
import math
from datetime import datetime, timedelta

from pptx import Presentation
from pptx.chart.data import (
    BubbleChartData, CategoryChartData, XyChartData,
)
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION,
)
from pptx.enum.dml import MSO_PATTERN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from specs import (
    BarSpec, ButterflySpec, ChartSpec, GanttSpec, LineSpec, MekkoSpec,
    PieSpec, ProcessSpec, ScatterSpec, TableSpec, UNIT_EXCEL, WaterfallSpec,
    decode_cell, spec_kind, weekend_ranges,
)
from theme import Theme

EMU_IN = 914400


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#"))


def _nice_max(v: float) -> float:
    """Round v up to a 'nice' axis maximum (1/2/2.5/5 × 10^k)."""
    if v <= 0:
        return 1.0
    exp = math.floor(math.log10(v))
    frac = v / (10 ** exp)
    for nice in (1.0, 1.2, 1.5, 2.0, 2.5, 3.0, 4.0, 5.0, 7.5, 10.0):
        if frac <= nice:
            return nice * (10 ** exp)
    return 10.0 * (10 ** exp)


# ---------------------------------------------------------------------- #
# Deck assembly
# ---------------------------------------------------------------------- #

def build_deck(
    specs: list[ChartSpec],
    theme: Theme,
    template: bytes | None = None,
    deck_title: str | None = None,
    *,
    groups: list[list[ChartSpec]] | None = None,
    agenda: str = "none",              # "none" | "front" | "chapters"
    harmonise_bars: bool = False,      # same value-axis max on stacked bars
) -> bytes:
    prs = Presentation(io.BytesIO(template)) if template else Presentation()
    if not template:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders))

    grouped = groups if groups is not None else [[s] for s in specs]
    grouped = [g for g in grouped if g]
    flat = [s for g in grouped for s in g]

    vmax_shared = None
    if harmonise_bars:
        cands = [s for s in flat
                 if isinstance(s, BarSpec) and s.mode == "stacked"
                 and s.orientation == "vertical" and not s.axis_break]
        if cands:
            vmax_shared = _nice_max(
                max(max(c.totals()) for c in cands) * 1.12)

    if deck_title:
        slide = prs.slides.add_slide(layout)
        _textbox(slide, deck_title, Inches(0.8), Inches(2.9),
                 prs.slide_width - Inches(1.6), Inches(1.1),
                 size=32, bold=True, color=theme.text, font=theme.font)
        _textbox(slide, datetime.now().strftime("%d %B %Y"),
                 Inches(0.8), Inches(4.0), Inches(5), Inches(0.4),
                 size=14, color=theme.muted, font=theme.font)

    section_titles = [g[0].title for g in grouped]
    if agenda == "front" and section_titles:
        _agenda_slide(prs, layout, section_titles, None, theme)

    for gi, group in enumerate(grouped):
        if agenda == "chapters" and section_titles:
            _agenda_slide(prs, layout, section_titles, gi, theme)
        slide = prs.slides.add_slide(layout)
        for spec, frame in zip(group, _grid_frames(prs, len(group))):
            _render_one(slide, spec, theme, frame,
                        small=len(group) > 1, vmax_shared=vmax_shared)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _grid_frames(prs, k: int) -> list[tuple]:
    """Full-slide, 2-up or 4-up chart regions (title drawn inside each)."""
    x0, y0 = Inches(0.5), Inches(0.25)
    w = prs.slide_width - Inches(1.0)
    h = prs.slide_height - Inches(0.70)
    gap = Inches(0.25)
    if k <= 1:
        return [(x0, y0, w, h)]
    if k == 2:
        cw = int((w - gap) / 2)
        return [(x0, y0, cw, h), (x0 + cw + gap, y0, cw, h)]
    cw, ch = int((w - gap) / 2), int((h - gap) / 2)
    cells = [(x0, y0), (x0 + cw + gap, y0),
             (x0, y0 + ch + gap), (x0 + cw + gap, y0 + ch + gap)]
    return [(cx, cy, cw, ch) for cx, cy in cells[:k]]


def _render_one(slide, spec: ChartSpec, theme: Theme, outer,
                small: bool = False, vmax_shared: float | None = None):
    x, y, cx, cy = outer
    title_h = Inches(0.45 if small else 0.65)
    _textbox(slide, spec.title, x, y, cx, title_h,
             size=13 if small else 20, bold=True, color=theme.text,
             font=theme.font)
    frame = (x, y + title_h, cx, cy - title_h - Inches(0.10))
    if isinstance(spec, BarSpec):
        _bar_slide(slide, spec, theme, frame, vmax_shared=vmax_shared)
    elif isinstance(spec, WaterfallSpec):
        _waterfall_slide(slide, spec, theme, frame)
    elif isinstance(spec, GanttSpec):
        _gantt_slide(slide, spec, theme, frame)
    elif isinstance(spec, MekkoSpec):
        _mekko_slide(slide, spec, theme, frame)
    elif isinstance(spec, TableSpec):
        _table_slide(slide, spec, theme, frame)
    elif isinstance(spec, LineSpec):
        _line_slide(slide, spec, theme, frame)
    elif isinstance(spec, PieSpec):
        _pie_slide(slide, spec, theme, frame)
    elif isinstance(spec, ButterflySpec):
        _butterfly_slide(slide, spec, theme, frame)
    elif isinstance(spec, ScatterSpec):
        _scatter_slide(slide, spec, theme, frame)
    elif isinstance(spec, ProcessSpec):
        _process_slide(slide, spec, theme, frame)
    else:
        raise TypeError(f"No pptx renderer for {spec_kind(spec)}")


def _agenda_slide(prs, layout, titles: list[str], current: int | None,
                  theme: Theme) -> None:
    slide = prs.slides.add_slide(layout)
    _textbox(slide, "Agenda", Inches(0.8), Inches(0.55), Inches(6),
             Inches(0.7), size=26, bold=True, color=theme.text,
             font=theme.font)
    y = Inches(1.7)
    for i, t in enumerate(titles):
        hot = current == i
        if hot:
            hl = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), y - Inches(0.06),
                prs.slide_width - Inches(1.3), Inches(0.44))
            hl.adjustments[0] = 0.5
            hl.fill.solid()
            hl.fill.fore_color.rgb = _rgb(theme.accent)
            hl.line.fill.background()
        _textbox(slide, f"{i + 1}", Inches(0.9), y, Inches(0.5),
                 Inches(0.35), size=14, bold=True,
                 color="#FFFFFF" if hot else theme.accent, font=theme.font)
        _textbox(slide, t, Inches(1.5), y,
                 prs.slide_width - Inches(2.3), Inches(0.35), size=14,
                 bold=hot, color="#FFFFFF" if hot else theme.text,
                 font=theme.font)
        y += Inches(0.52)


def _textbox(slide, text, x, y, cx, cy, *, size=12, bold=False,
             color="#1f2937", font="Arial", align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return tb


def _line(slide, x1, y1, x2, y2, color, width_pt=1.0, dashed=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    conn.line.color.rgb = _rgb(color)
    conn.line.width = Pt(width_pt)
    if dashed:
        # python-pptx has no dash API on connectors — write the XML directly.
        ln = conn.line._get_or_add_ln()
        dash = ln.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash",
            {"val": "dash"})
        ln.append(dash)
    return conn


# ---------------------------------------------------------------------- #
# Plot-area geometry model for native-chart overlays
# ---------------------------------------------------------------------- #

# Fractions of the chart graphic-frame occupied by non-plot furniture.
# Chosen to match PowerPoint's default auto-layout closely; the value axis
# max is pinned explicitly so the vertical mapping is exact up to these
# insets.
_INSET_L, _INSET_R, _INSET_T, _INSET_B = 0.075, 0.015, 0.035, 0.09
_LEGEND_B = 0.09  # extra bottom inset when a legend is shown


class _Geometry:
    def __init__(self, frame, n_cats: int, vmax: float, gap_width: int,
                 legend: bool):
        x, y, cx, cy = frame
        bottom = _INSET_B + (_LEGEND_B if legend else 0.0)
        self.plot_x = x + int(cx * _INSET_L)
        self.plot_w = int(cx * (1 - _INSET_L - _INSET_R))
        self.plot_y = y + int(cy * _INSET_T)
        self.plot_h = int(cy * (1 - _INSET_T - bottom))
        self.n = n_cats
        self.vmax = vmax
        # gap_width is the gap as % of bar width -> bar fraction of slot.
        self.bar_frac = 1.0 / (1.0 + gap_width / 100.0)

    def cat_x(self, i: float) -> int:
        return int(self.plot_x + self.plot_w * (i + 0.5) / self.n)

    def bar_halfwidth(self) -> int:
        return int(self.plot_w / self.n * self.bar_frac / 2)

    def y_of(self, v: float) -> int:
        return int(self.plot_y + self.plot_h * (1.0 - v / self.vmax))


# ---------------------------------------------------------------------- #
# Bars — native chart + overlays
# ---------------------------------------------------------------------- #

_XL_BY_MODE = {
    ("stacked", "vertical"): XL_CHART_TYPE.COLUMN_STACKED,
    ("clustered", "vertical"): XL_CHART_TYPE.COLUMN_CLUSTERED,
    ("stacked100", "vertical"): XL_CHART_TYPE.COLUMN_STACKED_100,
    ("stacked", "horizontal"): XL_CHART_TYPE.BAR_STACKED,
    ("clustered", "horizontal"): XL_CHART_TYPE.BAR_CLUSTERED,
    ("stacked100", "horizontal"): XL_CHART_TYPE.BAR_STACKED_100,
}
_GAP_WIDTH = 60


def _delete_axis(axis) -> None:
    el = axis._element
    d = el.find(qn("c:delete"))
    if d is None:
        d = el.makeelement(qn("c:delete"), {})
        el.find(qn("c:scaling")).addnext(d)
    d.set("val", "1")


def _chart_transparent(chart) -> None:
    """No chart-area fill/border so an overlay chart shows what's below."""
    cs = chart._chartSpace
    sp = cs.find(qn("c:spPr"))
    if sp is None:
        sp = cs.makeelement(qn("c:spPr"), {})
        cs.find(qn("c:chart")).addnext(sp)
    sp.append(sp.makeelement(qn("a:noFill"), {}))
    ln = sp.makeelement(qn("a:ln"), {})
    ln.append(ln.makeelement(qn("a:noFill"), {}))
    sp.append(ln)


def _manual_plot_layout(chart, xf: float, yf: float, wf: float,
                        hf: float) -> None:
    """Pin the inner plot area to exact fractions of the graphic frame —
    makes the overlay geometry exact and lets two charts align perfectly."""
    plot_area = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    layout = plot_area.find(qn("c:layout"))
    if layout is None:
        layout = plot_area.makeelement(qn("c:layout"), {})
        plot_area.insert(0, layout)
    for child in list(layout):
        layout.remove(child)
    man = layout.makeelement(qn("c:manualLayout"), {})
    layout.append(man)
    for tag, val in (("layoutTarget", "inner"), ("xMode", "edge"),
                     ("yMode", "edge"), ("x", str(xf)), ("y", str(yf)),
                     ("w", str(wf)), ("h", str(hf))):
        el = man.makeelement(qn("c:" + tag), {})
        el.set("val", val)
        man.append(el)


def _bar_slide(slide, spec: BarSpec, theme: Theme, frame,
               vmax_shared: float | None = None) -> None:
    if (spec.axis_break and spec.orientation == "vertical"
            and spec.mode in ("stacked", "clustered")):
        _bar_break_slide(slide, spec, theme, frame)
        return

    x, y, cx, cy = frame
    arrows = (spec.resolved_deltas()
              if spec.mode == "stacked" and spec.orientation == "vertical"
              else [])
    gutter_arrows = [a for a in arrows if not a[4].startswith("CAGR")]
    if gutter_arrows:
        gutter_w = Inches(0.9 + 0.75 * len(gutter_arrows))
        cx = cx - gutter_w

    nf = UNIT_EXCEL.get(spec.unit) or "#,##0"
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories
    for s in spec.series:
        chart_data.add_series(s.name, tuple(v or 0.0 for v in s.values))

    gf = slide.shapes.add_chart(
        _XL_BY_MODE[(spec.mode, spec.orientation)], x, y, cx, cy, chart_data)
    chart = gf.chart
    chart.has_title = False
    chart.font.name = theme.font
    chart.font.size = Pt(theme.base_font_pt)

    legend = len(spec.series) > 1
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.gap_width = _GAP_WIDTH
    if spec.mode in ("stacked", "stacked100"):
        plot.overlap = 100
    for k, s in enumerate(spec.series):
        chart.series[k].format.fill.solid()
        chart.series[k].format.fill.fore_color.rgb = _rgb(
            s.color or theme.color(k))

    if spec.show_values:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = nf
        dl.number_format_is_linked = False
        dl.font.size = Pt(10)
        dl.position = (XL_LABEL_POSITION.CENTER
                       if spec.mode in ("stacked", "stacked100")
                       else XL_LABEL_POSITION.OUTSIDE_END)
    if spec.unit:
        chart.value_axis.tick_labels.number_format = nf
        chart.value_axis.tick_labels.number_format_is_linked = False

    totals = spec.totals()
    combo = (spec.overlay_lines and spec.orientation == "vertical"
             and spec.mode in ("stacked", "clustered"))
    vmax = None
    if spec.mode == "stacked" and totals:
        vmax = vmax_shared or _nice_max(max(totals) * 1.12)
    elif combo:
        peak = max((v or 0.0) for s in spec.series for v in s.values)
        vmax = _nice_max(peak * 1.12)
    if combo and vmax:
        line_peak = max((v or 0.0) for s in spec.overlay_lines
                        for v in s.values)
        vmax = max(vmax, _nice_max(line_peak * 1.12))

    overlays = (spec.mode == "stacked" and spec.orientation == "vertical"
                and vmax)
    if overlays or combo:
        # Pin axis AND plot area so the geometry model is exact.
        chart.value_axis.minimum_scale = 0.0
        chart.value_axis.maximum_scale = vmax
        bottom = _INSET_B + (_LEGEND_B if legend else 0.0)
        _manual_plot_layout(chart, _INSET_L, _INSET_T,
                            1 - _INSET_L - _INSET_R,
                            1 - _INSET_T - bottom)
    if combo:
        _overlay_line_chart(slide, spec, theme, (x, y, cx, cy), vmax, legend)
    if overlays:
        geo = _Geometry(( x, y, cx, cy), len(spec.categories), vmax,
                        _GAP_WIDTH, legend)

        if spec.show_totals:
            for i, t in enumerate(totals):
                _textbox(slide, spec.fmt(t),
                         geo.cat_x(i) - Inches(0.6), geo.y_of(t) - Inches(0.3),
                         Inches(1.2), Inches(0.25), size=11, bold=True,
                         color=theme.text, font=theme.font,
                         align=PP_ALIGN.CENTER)

        gutter_x = x + cx + Inches(0.35)
        for (i, j, ti, tj, label) in arrows:
            if label.startswith("CAGR"):
                _cagr_arrow(slide, geo, i, j, ti, tj, label, theme)
                continue
            color = theme.positive if tj < ti else theme.negative
            for cat, tot in ((i, ti), (j, tj)):
                _line(slide, geo.cat_x(cat) + geo.bar_halfwidth(),
                      geo.y_of(tot), gutter_x + Inches(0.10), geo.y_of(tot),
                      theme.muted, 0.75, dashed=True)
            top, bot = sorted((geo.y_of(ti), geo.y_of(tj)))
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.UP_DOWN_ARROW, int(gutter_x), int(top),
                Inches(0.20), int(bot - top) or Inches(0.1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = _rgb(color)
            arrow.line.fill.background()
            _textbox(slide, label, gutter_x + Inches(0.26),
                     Emu(int((top + bot) / 2)) - Inches(0.12),
                     Inches(1.15), Inches(0.24), size=11, bold=True,
                     color=color, font=theme.font)
            gutter_x = gutter_x + Inches(0.75)


def _cagr_arrow(slide, geo, i, j, ti, tj, label, theme: Theme) -> None:
    lift = Inches(0.28)
    x1, y1 = geo.cat_x(i), geo.y_of(ti) - lift
    x2, y2 = geo.cat_x(j), geo.y_of(tj) - lift
    length = int(math.hypot(x2 - x1, y2 - y1))
    angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
    cxm, cym = (x1 + x2) // 2, (y1 + y2) // 2
    thick = Inches(0.16)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, cxm - length // 2, cym - int(thick) // 2,
        length, thick)
    arrow.rotation = angle
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = _rgb(theme.accent)
    arrow.line.fill.background()
    _textbox(slide, label, Emu(cxm) - Inches(0.9),
             Emu(min(y1, y2)) - Inches(0.35), Inches(1.8), Inches(0.25),
             size=11, bold=True, color=theme.accent, font=theme.font,
             align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Waterfall — native stacked column, invisible base
# ---------------------------------------------------------------------- #

def _waterfall_slide(slide, spec: WaterfallSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    steps = spec.resolved()
    base = [s.base if s.kind != "total" else 0.0 for s in steps]
    inc = [s.delta if s.kind == "increase" else 0.0 for s in steps]
    dec = [-s.delta if s.kind == "decrease" else 0.0 for s in steps]
    tot = [s.cumulative if s.kind == "total" else 0.0 for s in steps]

    chart_data = CategoryChartData()
    chart_data.categories = [s.label for s in steps]
    chart_data.add_series("_base", tuple(base))
    chart_data.add_series("Increase", tuple(inc))
    chart_data.add_series("Decrease", tuple(dec))
    chart_data.add_series("Total", tuple(tot))

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = False
    chart.font.name = theme.font
    chart.font.size = Pt(theme.base_font_pt)
    plot = chart.plots[0]
    plot.gap_width = 45
    plot.overlap = 100

    chart.series[0].format.fill.background()
    for k, color in ((1, theme.negative), (2, theme.positive),
                     (3, theme.total)):
        chart.series[k].format.fill.solid()
        chart.series[k].format.fill.fore_color.rgb = _rgb(color)

    peak = max((max(s.cumulative, s.cumulative - min(s.delta, 0.0))
                for s in steps), default=1.0)
    vmax = _nice_max(peak * 1.15)
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = vmax

    # Signed per-point labels (native labels would show magnitudes only).
    if spec.show_values:
        for idx, s in enumerate(steps):
            series_i = {"increase": 1, "decrease": 2, "total": 3}[s.kind]
            point = chart.series[series_i].points[idx]
            dl = point.data_label
            tf = dl.text_frame
            tf.text = (f"{s.cumulative:,.0f}" if s.kind == "total"
                       else spec.number_format.format(s.delta))
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(10)
            run.font.bold = s.kind == "total"
            run.font.name = theme.font
            if s.kind == "total":
                run.font.color.rgb = _rgb("#FFFFFF")

    if spec.show_connectors:
        geo = _Geometry((x, y, cx, cy), len(steps), vmax, 45, legend=False)
        for i in range(len(steps) - 1):
            level = steps[i].cumulative
            _line(slide,
                  geo.cat_x(i) + geo.bar_halfwidth(), geo.y_of(level),
                  geo.cat_x(i + 1) - geo.bar_halfwidth(), geo.y_of(level),
                  theme.muted, 0.75)


# ---------------------------------------------------------------------- #
# Gantt — deterministic shapes
# ---------------------------------------------------------------------- #

def _month_starts(t0: datetime, t1: datetime) -> list[datetime]:
    out = []
    ym = (t0.year, t0.month)
    while True:
        d = datetime(ym[0], ym[1], 1)
        if d > t1:
            break
        if d >= t0:
            out.append(d)
        ym = (ym[0] + 1, 1) if ym[1] == 12 else (ym[0], ym[1] + 1)
    return out


def _gantt_slide(slide, spec: GanttSpec, theme: Theme, frame) -> None:
    """think-cell feature set: calendar header (years+months), row banding,
    curtains, weekend shading, labelled date lines, phase brackets, bar
    styles (solid/striped/open), date & duration labels, milestone date
    labels, remark column, multiple bars per row."""
    x, y, cx, cy = frame
    span = spec.span()
    if span is None or not spec.items:
        _textbox(slide, "No dated items.", x, y, cx, Inches(0.4),
                 size=12, color=theme.muted, font=theme.font)
        return
    pad = max((span[1] - span[0]).days * 0.03, 8)
    t0 = span[0] - timedelta(days=pad)
    t1 = span[1] + timedelta(days=pad)
    total_s = (t1 - t0).total_seconds()

    rows = spec.rows()
    n = len(rows)
    row_of = {label: i for i, label in enumerate(rows)}
    groups = spec.groups() or [""]
    color_of = {g: theme.color(i) for i, g in enumerate(groups)}
    has_remarks = spec.show_remarks and any(it.remark for it in spec.items)
    n_bracket_lvls = (max((lvl for _, lvl in spec.bracket_levels()),
                          default=-1) + 1)

    label_w = Inches(2.5)
    remark_w = Inches(1.7) if has_remarks else 0
    header_h = Inches(0.40)                     # years + months rows
    bracket_h = Inches(0.34) * n_bracket_lvls
    legend_h = Inches(0.32) if spec.groups() else 0
    chart_x = x + label_w
    chart_w = cx - label_w - remark_w
    body_y = y + header_h + bracket_h
    body_h = cy - header_h - bracket_h - legend_h
    body_bottom = body_y + body_h

    def x_of(dt: datetime) -> int:
        return int(chart_x + chart_w * (dt - t0).total_seconds() / total_s)

    row_h = int(body_h / n)
    bar_h = min(int(row_h * 0.55), int(Inches(0.30)))

    # --- background: row banding, weekend shading, curtains --------------
    for i in range(1, n, 2):
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, int(x), int(body_y + i * row_h),
            int(cx), row_h)
        band.fill.solid()
        band.fill.fore_color.rgb = _rgb(theme.band)
        band.line.fill.background()
    if spec.weekend_shading:
        for (ws, we) in weekend_ranges(t0, t1):
            wrect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_of(ws), int(body_y),
                max(x_of(we) - x_of(ws), int(Inches(0.01))), int(body_h))
            wrect.fill.solid()
            wrect.fill.fore_color.rgb = _rgb(theme.grid)
            wrect.line.fill.background()
    for c in spec.curtains:
        if c.start is None or c.end is None:
            continue
        cur = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_of(c.start), int(body_y),
            max(x_of(c.end) - x_of(c.start), int(Inches(0.02))),
            int(body_h))
        cur.fill.solid()
        cur.fill.fore_color.rgb = _rgb(c.color or theme.curtain)
        cur.line.fill.background()
        if c.label:
            _textbox(slide, c.label,
                     Emu((x_of(c.start) + x_of(c.end)) // 2) - Inches(0.8),
                     body_y + Inches(0.03), Inches(1.6), Inches(0.2),
                     size=9, color=theme.muted, font=theme.font,
                     align=PP_ALIGN.CENTER)

    # --- calendar header: years row + months (or quarters) row -----------
    months = _month_starts(t0, t1)
    n_months = len(months)
    if n_months > 30:                    # quarter labels when span is long
        ticks = [m for m in months if m.month in (1, 4, 7, 10)]
        fmt = "Q%q"
    else:
        ticks, fmt = months, "%b"
    step = max(1, math.ceil(len(ticks) / 20))
    years_row_y = y
    months_row_y = y + Inches(0.19)
    for k, m in enumerate(ticks):
        _line(slide, x_of(m), int(months_row_y) + int(Inches(0.18)),
              x_of(m), body_bottom, theme.grid, 0.5)
        if k % step == 0:
            label = (f"Q{(m.month - 1) // 3 + 1}" if fmt == "Q%q"
                     else m.strftime("%b"))
            _textbox(slide, label, Emu(x_of(m)) - Inches(0.3),
                     months_row_y, Inches(0.6), Inches(0.18), size=8,
                     color=theme.muted, font=theme.font,
                     align=PP_ALIGN.CENTER)
    for yr in sorted({m.year for m in months}):
        y_start = max(datetime(yr, 1, 1), t0)
        y_end = min(datetime(yr + 1, 1, 1), t1)
        if (y_end - y_start).days < 20:
            continue
        _line(slide, x_of(y_start), int(years_row_y),
              x_of(y_start), int(years_row_y) + int(Inches(0.36)),
              theme.muted, 0.75)
        _textbox(slide, str(yr),
                 Emu((x_of(y_start) + x_of(y_end)) // 2) - Inches(0.4),
                 years_row_y, Inches(0.8), Inches(0.18), size=9, bold=True,
                 color=theme.text, font=theme.font, align=PP_ALIGN.CENTER)

    # --- phase brackets ---------------------------------------------------
    for b, lvl in spec.bracket_levels():
        by = int(y + header_h + Inches(0.34) * lvl + Inches(0.22))
        _line(slide, x_of(b.start), by, x_of(b.end), by, theme.text, 1.3)
        for bx in (b.start, b.end):
            _line(slide, x_of(bx), by, x_of(bx), by + int(Inches(0.07)),
                  theme.text, 1.3)
        if b.label:
            _textbox(slide, b.label,
                     Emu((x_of(b.start) + x_of(b.end)) // 2) - Inches(1.2),
                     Emu(by) - Inches(0.20), Inches(2.4), Inches(0.18),
                     size=9, bold=True, color=theme.text, font=theme.font,
                     align=PP_ALIGN.CENTER)

    # --- rows: labels, bars, milestones, per-bar labels, remarks ----------
    for label, i in row_of.items():
        mid_y = body_y + i * row_h + row_h // 2
        _textbox(slide, label, x, Emu(int(mid_y)) - Inches(0.10),
                 label_w - Inches(0.15), Inches(0.22), size=10,
                 color=theme.text, font=theme.font)

    if has_remarks:
        remark_by_row: dict[int, str] = {}
        for it in spec.items:
            if it.remark and row_of[it.label] not in remark_by_row:
                remark_by_row[row_of[it.label]] = it.remark
        for i, remark in remark_by_row.items():
            mid_y = body_y + i * row_h + row_h // 2
            _textbox(slide, remark, x + cx - remark_w + Inches(0.1),
                     Emu(int(mid_y)) - Inches(0.10), remark_w - Inches(0.1),
                     Inches(0.22), size=9, color=theme.muted,
                     font=theme.font)

    for it in spec.items:
        mid_y = body_y + row_of[it.label] * row_h + row_h // 2
        color = color_of.get(it.group or groups[0], theme.color(0))
        if it.kind == "milestone":
            when = it.start or it.finish
            if when is None:
                continue
            size = Inches(0.20)
            d = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND, x_of(when) - int(size) // 2,
                int(mid_y) - int(size) // 2, size, size)
            d.fill.solid()
            d.fill.fore_color.rgb = _rgb(color)
            d.line.color.rgb = _rgb("#FFFFFF")
            d.line.width = Pt(0.75)
            if spec.show_date_labels:
                _textbox(slide, f"{when:%d %b %y}",
                         Emu(x_of(when)) + Inches(0.14),
                         Emu(int(mid_y)) - Inches(0.09), Inches(0.9),
                         Inches(0.18), size=8, color=theme.muted,
                         font=theme.font)
        elif it.start and it.finish:
            x1, x2 = x_of(it.start), x_of(max(it.finish, it.start))
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x1,
                int(mid_y) - bar_h // 2,
                max(x2 - x1, int(Inches(0.03))), bar_h)
            bar.adjustments[0] = 0.5
            if it.style == "striped":
                bar.fill.patterned()
                bar.fill.pattern = MSO_PATTERN.LIGHT_DOWNWARD_DIAGONAL
                bar.fill.fore_color.rgb = _rgb(color)
                bar.fill.back_color.rgb = _rgb("#FFFFFF")
                bar.line.color.rgb = _rgb(color)
                bar.line.width = Pt(0.75)
            elif it.style == "open":
                bar.fill.background()
                bar.line.color.rgb = _rgb(color)
                bar.line.width = Pt(1.3)
            else:
                bar.fill.solid()
                bar.fill.fore_color.rgb = _rgb(color)
                bar.line.fill.background()
            if spec.show_date_labels:
                _textbox(slide, f"{it.start:%d %b %y}",
                         Emu(x1) - Inches(0.95),
                         Emu(int(mid_y)) - Inches(0.09), Inches(0.88),
                         Inches(0.18), size=8, color=theme.muted,
                         font=theme.font, align=PP_ALIGN.RIGHT)
                _textbox(slide, f"{it.finish:%d %b %y}",
                         Emu(x2) + Inches(0.06),
                         Emu(int(mid_y)) - Inches(0.09), Inches(0.88),
                         Inches(0.18), size=8, color=theme.muted,
                         font=theme.font)
            if spec.show_durations:
                days = (it.finish - it.start).days
                _textbox(slide, f"{days}d",
                         Emu((x1 + x2) // 2) - Inches(0.3),
                         Emu(int(mid_y) - bar_h // 2) - Inches(0.17),
                         Inches(0.6), Inches(0.16), size=8,
                         color=theme.text, font=theme.font,
                         align=PP_ALIGN.CENTER)

    # --- labelled date lines (incl. legacy today) -------------------------
    for dl in spec.all_date_lines():
        if dl.date is None or not (t0 <= dl.date <= t1):
            continue
        color = dl.color or theme.accent
        _line(slide, x_of(dl.date), int(body_y), x_of(dl.date),
              body_bottom, color, 1.25, dashed=True)
        if dl.label:
            _textbox(slide, dl.label, Emu(x_of(dl.date)) + Inches(0.04),
                     body_y + Inches(0.01), Inches(1.2), Inches(0.18),
                     size=9, color=color, font=theme.font)

    # --- legend ------------------------------------------------------------
    if spec.groups():
        lx = chart_x
        ly = body_bottom + Inches(0.06)
        for g in groups:
            chip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(lx), int(ly), Inches(0.16),
                Inches(0.16))
            chip.fill.solid()
            chip.fill.fore_color.rgb = _rgb(color_of[g])
            chip.line.fill.background()
            _textbox(slide, g, lx + Inches(0.22), Emu(int(ly)) - Inches(0.02),
                     Inches(1.6), Inches(0.22), size=10, color=theme.text,
                     font=theme.font)
            lx = lx + Inches(0.22) + Inches(0.13) * (len(g) + 4)


# ---------------------------------------------------------------------- #
# Marimekko — deterministic shapes
# ---------------------------------------------------------------------- #

def _mekko_slide(slide, spec: MekkoSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    header_h = Inches(0.55)
    legend_h = Inches(0.35)
    body_y = y + header_h
    body_h = cy - header_h - legend_h
    shares = spec.column_shares()
    totals = spec.column_totals()

    left = x
    for i, cat in enumerate(spec.categories):
        col_w = int(cx * shares[i])
        _textbox(slide, f"{cat}\n{spec.number_format.format(totals[i])}",
                 left, y, Emu(col_w), header_h, size=10, bold=True,
                 color=theme.text, font=theme.font, align=PP_ALIGN.CENTER)
        seg_y = body_y
        col_total = totals[i] or 1.0
        for k, s in enumerate(spec.series):
            frac = (s.values[i] or 0.0) / col_total
            seg_h = int(body_h * frac)
            if seg_h <= 0:
                continue
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(left), int(seg_y),
                max(col_w - int(Inches(0.02)), int(Inches(0.05))), seg_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = _rgb(s.color or theme.color(k))
            rect.line.color.rgb = _rgb("#FFFFFF")
            rect.line.width = Pt(1.0)
            if (spec.show_values and frac >= 0.07
                    and col_w >= int(Inches(0.75))):
                tb = _textbox(slide, f"{frac:.0%}", int(left),
                              Emu(int(seg_y + seg_h / 2)) - Inches(0.10),
                              Emu(col_w), Inches(0.2), size=10,
                              color="#FFFFFF", font=theme.font,
                              align=PP_ALIGN.CENTER)
                tb.text_frame.word_wrap = False
            seg_y += seg_h
        left += col_w

    lx = x
    ly = y + cy - legend_h + Inches(0.08)
    for k, s in enumerate(spec.series):
        chip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, int(lx), int(ly), Inches(0.16), Inches(0.16))
        chip.fill.solid()
        chip.fill.fore_color.rgb = _rgb(s.color or theme.color(k))
        chip.line.fill.background()
        _textbox(slide, s.name, lx + Inches(0.22),
                 Emu(int(ly)) - Inches(0.02), Inches(2.2), Inches(0.22),
                 size=10, color=theme.text, font=theme.font)
        lx = lx + Inches(0.22) + Inches(0.13) * (len(s.name) + 4)


# ---------------------------------------------------------------------- #
# Table — native
# ---------------------------------------------------------------------- #

def _table_slide(slide, spec: TableSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    rows, cols = len(spec.rows) + 1, len(spec.columns)
    height = min(cy, Inches(0.42) * rows)
    tbl = slide.shapes.add_table(rows, cols, x, y, cx, int(height)).table
    for c, name in enumerate(spec.columns):
        cell = tbl.cell(0, c)
        cell.text = str(name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme.accent)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(12)
        para.font.color.rgb = _rgb("#FFFFFF")
        para.font.name = theme.font
    semantic = {"accent": theme.accent, "positive": theme.positive,
                "negative": theme.negative, "warning": theme.warning,
                "muted": theme.muted}
    for r, row in enumerate(spec.rows, start=1):
        for c in range(cols):
            cell = tbl.cell(r, c)
            raw = str(row[c]) if c < len(row) else ""
            text, key = decode_cell(raw)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(
                "#FFFFFF" if r % 2 else "#F4F6F8")
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14 if key else 11.5)
            para.font.bold = bool(key)
            para.font.color.rgb = _rgb(semantic.get(key) or theme.text)
            para.font.name = theme.font
            if key:
                para.alignment = PP_ALIGN.CENTER


# ---------------------------------------------------------------------- #
# Combo overlay: transparent line chart pinned over the bar chart
# ---------------------------------------------------------------------- #

def _overlay_line_chart(slide, spec: BarSpec, theme: Theme, frame,
                        vmax: float, legend: bool) -> None:
    x, y, cx, cy = frame
    cd = CategoryChartData()
    cd.categories = spec.categories
    for s in spec.overlay_lines:
        cd.add_series(s.name, tuple(v or 0.0 for v in s.values))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    _chart_transparent(ch)
    ch.value_axis.minimum_scale = 0.0
    ch.value_axis.maximum_scale = vmax
    ch.value_axis.has_major_gridlines = False
    _delete_axis(ch.value_axis)
    _delete_axis(ch.category_axis)
    bottom = _INSET_B + (_LEGEND_B if legend else 0.0)
    _manual_plot_layout(ch, _INSET_L, _INSET_T, 1 - _INSET_L - _INSET_R,
                        1 - _INSET_T - bottom)
    geo = _Geometry((x, y, cx, cy), len(spec.categories), vmax,
                    _GAP_WIDTH, legend)
    for k, s in enumerate(spec.overlay_lines):
        color = s.color or theme.color(len(spec.series) + k)
        ser = ch.series[k]
        ser.format.line.color.rgb = _rgb(color)
        ser.format.line.width = Pt(2.25)
        ser.smooth = False
        last = max((i for i, v in enumerate(s.values) if v is not None),
                   default=None)
        if last is not None:
            _textbox(slide, s.name,
                     geo.cat_x(last) + Inches(0.12),
                     Emu(geo.y_of(s.values[last] or 0.0)) - Inches(0.12),
                     Inches(1.4), Inches(0.22), size=10, bold=True,
                     color=color, font=theme.font)


# ---------------------------------------------------------------------- #
# Bar chart with a value-axis break — deterministic shapes
# ---------------------------------------------------------------------- #

def _bar_break_slide(slide, spec: BarSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    n = len(spec.categories)
    stacked = spec.mode == "stacked"
    totals = spec.totals()
    vmax = (max(totals) if stacked
            else max((v or 0.0) for s in spec.series for v in s.values)
            ) * 1.05
    t, ticks = spec.axis_break.transform(vmax)

    axis_w = Inches(0.6)
    cat_h = Inches(0.30)
    legend_h = Inches(0.30) if len(spec.series) > 1 else 0
    plot_x = x + axis_w
    plot_w = cx - axis_w
    plot_y = y + Inches(0.15)
    plot_h = cy - Inches(0.15) - cat_h - legend_h
    base_y = plot_y + plot_h

    def y_of(display: float) -> int:
        return int(base_y - plot_h * display)

    for pos, val in ticks:
        _line(slide, int(plot_x), y_of(pos), int(x + cx), y_of(pos),
              theme.grid, 0.5)
        _textbox(slide, spec.fmt(val), x, Emu(y_of(pos)) - Inches(0.09),
                 axis_w - Inches(0.08), Inches(0.18), size=9,
                 color=theme.muted, font=theme.font, align=PP_ALIGN.RIGHT)

    slot = plot_w / n
    for i, cat in enumerate(spec.categories):
        cx_center = int(plot_x + slot * (i + 0.5))
        _textbox(slide, cat, Emu(cx_center) - Inches(0.7),
                 Emu(int(base_y)) + Inches(0.05), Inches(1.4), Inches(0.2),
                 size=10, color=theme.text, font=theme.font,
                 align=PP_ALIGN.CENTER)
        if stacked:
            bw = int(slot * 0.62)
            cum = 0.0
            for k, s in enumerate(spec.series):
                v = s.values[i] or 0.0
                h = int(plot_h * (t(cum + v) - t(cum)))
                top = y_of(t(cum + v))
                if h > 0:
                    seg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, cx_center - bw // 2, top, bw, h)
                    seg.fill.solid()
                    seg.fill.fore_color.rgb = _rgb(s.color or theme.color(k))
                    seg.line.fill.background()
                    if spec.show_values and h > int(Inches(0.22)):
                        _textbox(slide, spec.fmt(v),
                                 Emu(cx_center) - Inches(0.55),
                                 Emu(top + h // 2) - Inches(0.09),
                                 Inches(1.1), Inches(0.18), size=9,
                                 color="#FFFFFF", font=theme.font,
                                 align=PP_ALIGN.CENTER)
                cum += v
            if spec.show_totals:
                _textbox(slide, spec.fmt(totals[i]),
                         Emu(cx_center) - Inches(0.55),
                         Emu(y_of(t(totals[i]))) - Inches(0.28),
                         Inches(1.1), Inches(0.2), size=10, bold=True,
                         color=theme.text, font=theme.font,
                         align=PP_ALIGN.CENTER)
        else:
            k_n = len(spec.series)
            bw = int(slot * 0.62 / k_n)
            for k, s in enumerate(spec.series):
                v = s.values[i] or 0.0
                h = int(plot_h * t(v))
                left = cx_center - int(slot * 0.31) + k * bw
                if h > 0:
                    seg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, left, y_of(t(v)), bw, h)
                    seg.fill.solid()
                    seg.fill.fore_color.rgb = _rgb(s.color or theme.color(k))
                    seg.line.fill.background()
                if spec.show_values:
                    _textbox(slide, spec.fmt(v), Emu(left) - Inches(0.3),
                             Emu(y_of(t(v))) - Inches(0.22),
                             Emu(bw) + Inches(0.6), Inches(0.18), size=8,
                             color=theme.text, font=theme.font,
                             align=PP_ALIGN.CENTER)

    # Squiggle: slanted white band with grey edges across the gap.
    mid = (t(spec.axis_break.start) + t(spec.axis_break.end)) / 2
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(x + axis_w - Inches(0.1)),
        y_of(mid) - int(Inches(0.055)),
        int(plot_w + Inches(0.1)), int(Inches(0.11)))
    band.rotation = -3
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb("#FFFFFF")
    band.line.fill.background()
    for off in (-Inches(0.055), Inches(0.055)):
        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, int(x + axis_w - Inches(0.1)),
            y_of(mid) + int(off) - int(Pt(0.6)),
            int(plot_w + Inches(0.1)), int(Pt(1.2)))
        edge.rotation = -3
        edge.fill.solid()
        edge.fill.fore_color.rgb = _rgb(theme.muted)
        edge.line.fill.background()

    if len(spec.series) > 1:
        lx = plot_x
        ly = base_y + cat_h + Inches(0.04)
        for k, s in enumerate(spec.series):
            chip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(lx), int(ly), Inches(0.16),
                Inches(0.16))
            chip.fill.solid()
            chip.fill.fore_color.rgb = _rgb(s.color or theme.color(k))
            chip.line.fill.background()
            _textbox(slide, s.name, lx + Inches(0.22),
                     Emu(int(ly)) - Inches(0.02), Inches(1.8), Inches(0.22),
                     size=10, color=theme.text, font=theme.font)
            lx = lx + Inches(0.22) + Inches(0.13) * (len(s.name) + 4)


# ---------------------------------------------------------------------- #
# Line / area — native chart
# ---------------------------------------------------------------------- #

def _line_slide(slide, spec: LineSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    cd = CategoryChartData()
    cd.categories = spec.categories
    for s in spec.series:
        cd.add_series(s.name, tuple(v if v is not None else 0.0
                                    for v in s.values))
    xl = (XL_CHART_TYPE.AREA_STACKED if spec.mode == "area"
          else XL_CHART_TYPE.LINE_MARKERS)
    gf = slide.shapes.add_chart(xl, x, y, cx, cy, cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.name = theme.font
    ch.font.size = Pt(theme.base_font_pt)
    ch.has_legend = len(spec.series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    for k, s in enumerate(spec.series):
        color = _rgb(s.color or theme.color(k))
        ser = ch.series[k]
        if spec.mode == "area":
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = color
        else:
            ser.format.line.color.rgb = color
            ser.format.line.width = Pt(2.25)
            ser.smooth = False
    nf = UNIT_EXCEL.get(spec.unit) or "#,##0"
    if spec.show_values:
        plot = ch.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = nf
        dl.number_format_is_linked = False
        dl.font.size = Pt(9)
        if spec.mode == "line":
            dl.position = XL_LABEL_POSITION.ABOVE
    if spec.unit:
        ch.value_axis.tick_labels.number_format = nf
        ch.value_axis.tick_labels.number_format_is_linked = False


# ---------------------------------------------------------------------- #
# Pie / doughnut — native chart
# ---------------------------------------------------------------------- #

def _pie_slide(slide, spec: PieSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    cd = CategoryChartData()
    cd.categories = spec.labels
    cd.add_series("", tuple(spec.values))
    xl = XL_CHART_TYPE.DOUGHNUT if spec.doughnut else XL_CHART_TYPE.PIE
    gf = slide.shapes.add_chart(xl, x, y, cx, cy, cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.name = theme.font
    ch.font.size = Pt(theme.base_font_pt)
    for i in range(len(spec.labels)):
        pt = ch.series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _rgb(theme.color(i))
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = spec.show_pcts
    dl.show_value = not spec.show_pcts
    dl.font.size = Pt(11)
    if spec.doughnut:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        dl.show_category_name = False
    else:
        ch.has_legend = False
        dl.show_category_name = True
        dl.position = XL_LABEL_POSITION.OUTSIDE_END


# ---------------------------------------------------------------------- #
# Butterfly / tornado — native clustered bars, negatives shown absolute
# ---------------------------------------------------------------------- #

def _butterfly_slide(slide, spec: ButterflySpec, theme: Theme,
                     frame) -> None:
    x, y, cx, cy = frame
    cd = CategoryChartData()
    cd.categories = spec.categories
    cd.add_series(spec.left.name,
                  tuple(-(v or 0.0) for v in spec.left.values))
    cd.add_series(spec.right.name,
                  tuple(v or 0.0 for v in spec.right.values))
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.name = theme.font
    ch.font.size = Pt(theme.base_font_pt)
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    plot = ch.plots[0]
    plot.gap_width = 40
    plot.overlap = 100
    for k, s in enumerate((spec.left, spec.right)):
        ch.series[k].format.fill.solid()
        ch.series[k].format.fill.fore_color.rgb = _rgb(
            s.color or theme.color(k))
    base = UNIT_EXCEL.get(spec.unit) or "#,##0"
    nf_abs = f"{base};{base}"          # negatives displayed absolute
    if spec.show_values:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = nf_abs
        dl.number_format_is_linked = False
        dl.font.size = Pt(10)
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    ch.value_axis.tick_labels.number_format = nf_abs
    ch.value_axis.tick_labels.number_format_is_linked = False
    span = max([abs(v or 0.0) for v in
                spec.left.values + spec.right.values] or [1.0]) * 1.25
    ch.value_axis.minimum_scale = -span
    ch.value_axis.maximum_scale = span


# ---------------------------------------------------------------------- #
# Scatter / bubble — native XY chart
# ---------------------------------------------------------------------- #

def _scatter_slide(slide, spec: ScatterSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    groups = spec.groups() or [""]
    bubble = spec.bubble
    cd = BubbleChartData() if bubble else XyChartData()
    for g in groups:
        ser = cd.add_series(g or "Series 1")
        for p in spec.points:
            if (p.group or groups[0]) != g:
                continue
            if bubble:
                ser.add_data_point(p.x, p.y, max(p.size or 1.0, 0.1))
            else:
                ser.add_data_point(p.x, p.y)
    xl = XL_CHART_TYPE.BUBBLE if bubble else XL_CHART_TYPE.XY_SCATTER
    gf = slide.shapes.add_chart(xl, x, y, cx, cy, cd)
    ch = gf.chart
    ch.has_title = False
    ch.font.name = theme.font
    ch.font.size = Pt(theme.base_font_pt)
    ch.has_legend = bool(spec.groups())
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    for k, g in enumerate(groups):
        ser = ch.series[k]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = _rgb(theme.color(k))
        ser.format.line.fill.background()

    xs = [p.x for p in spec.points] or [0.0]
    ys = [p.y for p in spec.points] or [0.0]
    x_lo, x_hi = min(min(xs), 0.0), _nice_max(max(xs) * 1.15)
    y_lo, y_hi = min(min(ys), 0.0), _nice_max(max(ys) * 1.15)
    ch.category_axis.minimum_scale = x_lo
    ch.category_axis.maximum_scale = x_hi
    ch.value_axis.minimum_scale = y_lo
    ch.value_axis.maximum_scale = y_hi
    legend = ch.has_legend
    bottom = _INSET_B + (_LEGEND_B if legend else 0.0)
    _manual_plot_layout(ch, _INSET_L, _INSET_T, 1 - _INSET_L - _INSET_R,
                        1 - _INSET_T - bottom)

    # Geometry now exact → point labels + quadrant lines as shapes.
    plot_x = x + int(cx * _INSET_L)
    plot_w = int(cx * (1 - _INSET_L - _INSET_R))
    plot_y = y + int(cy * _INSET_T)
    plot_h = int(cy * (1 - _INSET_T - bottom))

    def px(vx: float) -> int:
        return int(plot_x + plot_w * (vx - x_lo) / (x_hi - x_lo))

    def py(vy: float) -> int:
        return int(plot_y + plot_h * (1 - (vy - y_lo) / (y_hi - y_lo)))

    for p in spec.points:
        _textbox(slide, p.label, Emu(px(p.x)) - Inches(0.7),
                 Emu(py(p.y)) - Inches(0.30), Inches(1.4), Inches(0.18),
                 size=9, color=theme.muted, font=theme.font,
                 align=PP_ALIGN.CENTER)
    if spec.quadrants:
        qx, qy = spec.quadrants
        if x_lo <= qx <= x_hi:
            _line(slide, px(qx), plot_y, px(qx), plot_y + plot_h,
                  theme.muted, 1.0, dashed=True)
        if y_lo <= qy <= y_hi:
            _line(slide, plot_x, py(qy), plot_x + plot_w, py(qy),
                  theme.muted, 1.0, dashed=True)
    if spec.x_title:
        _textbox(slide, spec.x_title, Emu(plot_x), y + cy - Inches(0.22),
                 Emu(plot_w), Inches(0.2), size=10, color=theme.muted,
                 font=theme.font, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Process flow — chevron strip
# ---------------------------------------------------------------------- #

def _process_slide(slide, spec: ProcessSpec, theme: Theme, frame) -> None:
    x, y, cx, cy = frame
    n = max(len(spec.steps), 1)
    h = min(int(Inches(1.15)), int(cy * 0.5))
    top = y + (cy - h) // 3
    gap = Inches(0.08)
    w = int((cx - gap * (n - 1)) / n)
    for i, step in enumerate(spec.steps):
        left = x + i * (w + gap)
        shape_type = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        sh = slide.shapes.add_shape(shape_type, int(left), int(top), w, h)
        hot = spec.highlight == i
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(theme.accent if hot else "#e4e8ec")
        sh.line.fill.background()
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = step
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = theme.font
        p.font.color.rgb = _rgb("#FFFFFF" if hot else theme.text)
